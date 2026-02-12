# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0

import io
import json
import logging
import os
from datetime import UTC, datetime

import boto3
from pptx import Presentation
from pptx.util import Pt

logger = logging.getLogger()
logger.setLevel(logging.INFO)


def create_presentation(title_, subtitle_, slides_data_, template_key_=None):
    try:
        s3 = boto3.client("s3")
        bucket_name = os.environ["S3_BUCKET"]

        # Get custom content from event
        title_text = title_.upper()
        subtitle_text = subtitle_.upper()
        slides_data = slides_data_
        template_key = template_key_ or "templates/template.pptx"

        # Download template from S3
        try:
            template_response = s3.get_object(Bucket=bucket_name, Key=template_key)
            template_data = template_response["Body"].read()

            # Load template
            template_buffer = io.BytesIO(template_data)
            prs = Presentation(template_buffer)

            # Update first slide with custom content
            if len(prs.slides) > 0:
                first_slide = prs.slides[0]

                # Update title
                if first_slide.shapes.title:
                    title_shape = first_slide.shapes.title
                    title_shape.text = title_text
                    # Format title
                    for paragraph in title_shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.bold = True
                            run.font.size = Pt(44)

                # Update subtitle
                if len(first_slide.placeholders) > 1:
                    subtitle_shape = first_slide.placeholders[1]
                    # Clear existing text frame
                    subtitle_shape.text_frame.clear()
                    # Add new paragraph without bullets
                    p = subtitle_shape.text_frame.paragraphs[0]
                    p.text = subtitle_text
                    p.level = 0
                    # Format subtitle
                    for run in p.runs:
                        run.font.bold = True
                        run.font.size = Pt(32)

        except s3.exceptions.NoSuchKey:
            # Create new presentation if template doesn't exist
            prs = Presentation()

            # Add title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)

            # Set and format title
            slide.shapes.title.text = title_text
            for paragraph in slide.shapes.title.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.size = Pt(44)

            # Set and format subtitle
            subtitle_placeholder = slide.placeholders[1]
            subtitle_placeholder.text_frame.clear()
            p = subtitle_placeholder.text_frame.paragraphs[0]
            p.text = subtitle_text
            p.level = 0
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(32)

        # Add additional slides from slides array
        total_slides_added = 0
        for slide_data in slides_data:
            slide_title = slide_data.get("title", "Slide Title").upper()
            slide_content = slide_data.get("content", "Slide Content")

            # Split content into lines
            lines = slide_content.split("\n")

            # Estimate lines per slide (assuming ~8-10 lines max for 14pt font)
            max_lines_per_slide = 8

            # Split lines into chunks if needed
            line_chunks = []
            for i in range(0, len(lines), max_lines_per_slide):
                chunk = lines[i : i + max_lines_per_slide]
                line_chunks.append(chunk)

            # Create slides for each chunk
            for chunk_index, line_chunk in enumerate(line_chunks):
                # Use content slide layout (layout 1)
                content_layout = prs.slide_layouts[1]
                new_slide = prs.slides.add_slide(content_layout)
                total_slides_added += 1

                # Set slide title (add part number if multiple slides)
                if len(line_chunks) > 1:
                    title_with_part = f"{slide_title} (Part {chunk_index + 1})"
                else:
                    title_with_part = slide_title

                new_slide.shapes.title.text = title_with_part
                for paragraph in new_slide.shapes.title.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                        run.font.size = Pt(36)

                # Set content with uniform 14pt font
                if len(new_slide.placeholders) > 1:
                    content_placeholder = new_slide.placeholders[1]

                    # Clear existing content and set new text
                    text_frame = content_placeholder.text_frame
                    text_frame.clear()
                    text_frame.word_wrap = True

                    # Add content line by line
                    for i, line in enumerate(line_chunk):
                        if not line.strip():  # Skip empty lines
                            continue

                        if i == 0:
                            # First paragraph already exists
                            p = text_frame.paragraphs[0]
                        else:
                            # Add new paragraphs for subsequent lines
                            p = text_frame.add_paragraph()

                        p.text = line.strip()

                        # Apply 14pt font size to all runs in the paragraph
                        for run in p.runs:
                            run.font.size = Pt(14)

        # Save to S3
        ppt_buffer = io.BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)

        key = f"presentations/presentation_{datetime.now(UTC).strftime('%Y%m%d_%H%M%S')}.pptx"

        s3.put_object(
            Bucket=bucket_name,
            Key=key,
            Body=ppt_buffer.getvalue(),
            ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

        # Generate CloudFront URL
        cloudfront_domain = os.environ["CLOUDFRONT_DOMAIN"]
        download_url = f"https://{cloudfront_domain}/{key}"

        return {
            "statusCode": 200,
            "body": json.dumps(
                {
                    "message": "PowerPoint created from template",
                    "key": key,
                    "title": title_text,
                    "subtitle": subtitle_text,
                    "slides_added": total_slides_added,
                    "download_url": download_url,
                }
            ),
        }
    except Exception as e:
        logger.error(f"Error creating presentation: {str(e)}")
        return {"statusCode": 500, "body": json.dumps({"error": str(e)})}


def lambda_handler(event, context):
    logger.info(f"INBOUND EVENT: {json.dumps(event)}")

    path = event.get("path", "")
    method = event.get("httpMethod", "")

    if method == "POST":
        try:
            body = json.loads(event.get("body", "{}"))
            method_name = body.get("method")
            request_id = body.get("id")

            if method_name == "initialize":
                response = {
                    "statusCode": 200,
                    "headers": {"Content-Type": "application/json"},
                    "body": json.dumps(
                        {
                            "jsonrpc": "2.0",
                            "id": request_id,
                            "result": {
                                "protocolVersion": "2025-03-26",
                                "capabilities": {"tools": {}, "logging": {}},
                                "serverInfo": {
                                    "name": "aws-mcp-server",
                                    "version": "1.0.0",
                                },
                            },
                        }
                    ),
                }
                logger.info(f"RESPONSE: {json.dumps(response)}")
                return response

            elif method_name == "notifications/initialized":
                response = {
                    "statusCode": 200,
                    "headers": {"Content-Type": "application/json"},
                    "body": "",
                }
                logger.info(f"RESPONSE: {json.dumps(response)}")
                return response

            elif method_name == "tools/list":
                response = {
                    "statusCode": 200,
                    "headers": {"Content-Type": "application/json"},
                    "body": json.dumps(
                        {
                            "jsonrpc": "2.0",
                            "id": request_id,
                            "result": {
                                "tools": [
                                    {
                                        "name": "ppt_creator",
                                        "description": "Create PowerPoint presentations with multiple slides",
                                        "inputSchema": {
                                            "type": "object",
                                            "properties": {
                                                "title": {
                                                    "type": "string",
                                                    "description": "Title of the presentation",
                                                },
                                                "subtitle": {
                                                    "type": "string",
                                                    "description": "Subtitle of the presentation",
                                                },
                                                "slides": {
                                                    "type": "array",
                                                    "description": "Array of slide objects",
                                                    "items": {
                                                        "type": "object",
                                                        "properties": {
                                                            "title": {
                                                                "type": "string",
                                                                "description": "Title of the slide",
                                                            },
                                                            "content": {
                                                                "type": "string",
                                                                "description": "Content of the slide",
                                                            },
                                                        },
                                                        "required": [
                                                            "title",
                                                            "content",
                                                        ],
                                                    },
                                                },
                                                "template_key": {
                                                    "type": "string",
                                                    "description": "S3 key path to the template file (e.g., templates/template.pptx). Defaults to templates/template.pptx if not provided",
                                                },
                                            },
                                            "required": ["title", "subtitle", "slides"],
                                        },
                                    }
                                ]
                            },
                        }
                    ),
                }
                logger.info(f"RESPONSE: {json.dumps(response)}")
                return response

            elif method_name == "tools/call":
                params = body.get("params", {})
                tool_name = params.get("name")
                arguments = params.get("arguments", {})

                if tool_name == "ppt_creator":
                    result = create_presentation(
                        arguments.get("title"),
                        arguments.get("subtitle"),
                        arguments.get("slides"),
                        arguments.get("template_key"),
                    )

                    if result.get("statusCode") == 200:
                        ppt_result = json.loads(result["body"])
                        response_body = json.dumps(
                            {
                                "jsonrpc": "2.0",
                                "id": request_id,
                                "result": {
                                    "content": [
                                        {
                                            "type": "text",
                                            "text": f"PowerPoint presentation '{ppt_result['title']}' created successfully!\n\nDownload URL: {ppt_result['download_url']}\n\nSlides added: {ppt_result['slides_added']}",
                                        }
                                    ]
                                },
                            }
                        )
                    else:
                        response_body = json.dumps(
                            {
                                "jsonrpc": "2.0",
                                "id": request_id,
                                "error": {
                                    "code": -32603,
                                    "message": f"Internal error: {json.loads(result['body'])['error']}",
                                },
                            }
                        )

                    response = {
                        "statusCode": 200,
                        "headers": {"Content-Type": "application/json"},
                        "body": response_body,
                    }

                    logger.info(f"RESPONSE: {json.dumps(response)}")
                    return response

                else:
                    response = {
                        "statusCode": 200,
                        "headers": {"Content-Type": "application/json"},
                        "body": json.dumps(
                            {
                                "jsonrpc": "2.0",
                                "id": request_id,
                                "error": {
                                    "code": -32601,
                                    "message": f"Method not found: {tool_name}",
                                },
                            }
                        ),
                    }
                    logger.info(f"RESPONSE: {json.dumps(response)}")
                    return response

            else:
                response = {
                    "statusCode": 200,
                    "headers": {"Content-Type": "application/json"},
                    "body": json.dumps(
                        {
                            "jsonrpc": "2.0",
                            "id": request_id,
                            "error": {
                                "code": -32601,
                                "message": f"Method not found: {method_name}",
                            },
                        }
                    ),
                }
                logger.info(f"RESPONSE: {json.dumps(response)}")
                return response

        except Exception as e:
            logger.error(f"JSON-RPC error: {str(e)}")
            response = {
                "statusCode": 200,
                "headers": {"Content-Type": "application/json"},
                "body": json.dumps(
                    {
                        "jsonrpc": "2.0",
                        "id": body.get("id") if "body" in locals() else None,
                        "error": {"code": -32700, "message": f"Parse error: {str(e)}"},
                    }
                ),
            }
            logger.info(f"RESPONSE: {json.dumps(response)}")
            return response

    # Handle non-POST requests
    return {"statusCode": 405, "body": json.dumps({"error": "Method not allowed"})}
